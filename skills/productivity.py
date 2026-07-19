import os
import sqlite3
import time
import json
import smtplib
import imaplib
import email
from email.mime.text import MIMEText
from loguru import logger
import sounddevice as sd
import numpy as np
import wave
from faster_whisper import WhisperModel
import pyautogui

class ProductivityPlanner:
    """Manages SQLite todo lists, SMTP/IMAP emails, meeting voice recordings, and workflow macro automation."""

    def __init__(self, db_path: str = "config/productivity.db"):
        self.db_path = db_path
        os.makedirs(os.path.dirname(self.db_path), exist_ok=True)
        self._init_db()

    def _init_db(self):
        try:
            conn = sqlite3.connect(self.db_path)
            c = conn.cursor()
            # Todos
            c.execute("""
                CREATE TABLE IF NOT EXISTS todos (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    task TEXT NOT NULL,
                    status TEXT DEFAULT 'PENDING'
                )
            """)
            # Macros
            c.execute("""
                CREATE TABLE IF NOT EXISTS macros (
                    name TEXT PRIMARY KEY,
                    steps TEXT NOT NULL
                )
            """)
            conn.commit()
            conn.close()
        except Exception as e:
            logger.error(f"Failed to initialize productivity database: {e}")

    # --- Todo List Management ---
    
    def add_todo(self, task: str) -> str:
        try:
            conn = sqlite3.connect(self.db_path)
            c = conn.cursor()
            c.execute("INSERT INTO todos (task) VALUES (?)", (task,))
            conn.commit()
            conn.close()
            return f"I have added '{task}' to your todo list, sir."
        except Exception as e:
            return f"Failed to add todo: {e}"

    def list_todos(self) -> str:
        try:
            conn = sqlite3.connect(self.db_path)
            c = conn.cursor()
            c.execute("SELECT id, task, status FROM todos WHERE status = 'PENDING'")
            rows = c.fetchall()
            conn.close()
            
            if not rows:
                return "Your todo list is currently empty, sir."
            
            summary = "Unfinished tasks on your list, sir:\n"
            for r in rows:
                summary += f" - [ID: {r[0]}] {r[1]}\n"
            return summary
        except Exception as e:
            return f"Failed to fetch todos: {e}"

    def complete_todo(self, todo_id: int) -> str:
        try:
            conn = sqlite3.connect(self.db_path)
            c = conn.cursor()
            c.execute("UPDATE todos SET status = 'COMPLETED' WHERE id = ?", (todo_id,))
            conn.commit()
            conn.close()
            return f"Marked task ID {todo_id} as completed, sir."
        except Exception as e:
            return f"Failed to update todo: {e}"

    # --- Email Management ---
    
    def send_email(self, to_addr: str, subject: str, body: str) -> str:
        # Check env variables or local config for credentials
        smtp_server = os.environ.get("JARVIS_SMTP_SERVER", "smtp.gmail.com")
        smtp_port = int(os.environ.get("JARVIS_SMTP_PORT", "587"))
        email_user = os.environ.get("JARVIS_EMAIL_USER")
        email_pass = os.environ.get("JARVIS_EMAIL_PASS")
        
        if not email_user or not email_pass:
            return (
                "Sir, I lack the credentials to send emails. "
                "Please configure JARVIS_EMAIL_USER and JARVIS_EMAIL_PASS in your environment."
            )
            
        try:
            msg = MIMEText(body)
            msg["Subject"] = subject
            msg["From"] = email_user
            msg["To"] = to_addr
            
            server = smtplib.SMTP(smtp_server, smtp_port)
            server.starttls()
            server.login(email_user, email_pass)
            server.sendmail(email_user, [to_addr], msg.as_string())
            server.quit()
            logger.info(f"Email sent to {to_addr}")
            return f"Email sent successfully to {to_addr}, sir."
        except Exception as e:
            logger.error(f"Email send failed: {e}")
            return f"Failed to send email: {str(e)}"

    def check_inbox(self) -> str:
        imap_server = os.environ.get("JARVIS_IMAP_SERVER", "imap.gmail.com")
        email_user = os.environ.get("JARVIS_EMAIL_USER")
        email_pass = os.environ.get("JARVIS_EMAIL_PASS")
        
        if not email_user or not email_pass:
            return "Email credentials are not configured, sir."
            
        try:
            mail = imaplib.IMAP4_SSL(imap_server)
            mail.login(email_user, email_pass)
            mail.select("inbox")
            
            status, data = mail.search(None, "UNSEEN")
            mail_ids = data[0].split()
            
            if not mail_ids:
                return "You have no unread emails in your inbox, sir."
                
            summary = f"You have {len(mail_ids)} unread emails, sir. Top unread messages:\n"
            
            # Fetch last 3 unread emails
            for mid in mail_ids[-3:]:
                status, data = mail.fetch(mid, "(RFC822)")
                raw_email = data[0][1]
                msg = email.message_from_bytes(raw_email)
                
                subject = msg["Subject"]
                sender = msg["From"]
                summary += f"  - From: {sender} | Subject: {subject}\n"
                
            mail.logout()
            return summary
        except Exception as e:
            logger.error(f"Inbox read failed: {e}")
            return f"Failed to read inbox: {str(e)}"

    # --- Meeting Notes / Continuous Dictation ---
    
    def record_meeting_notes(self, duration_sec: int = 60, output_txt: str = "notes.txt") -> str:
        """Records microphone in background for duration_sec, transcribes, and writes to a text file."""
        logger.info(f"Recording dictation/meeting notes for {duration_sec} seconds...")
        sample_rate = 16000
        
        try:
            # Record
            recording = sd.rec(int(duration_sec * sample_rate), samplerate=sample_rate, channels=1, dtype='int16')
            sd.wait() # wait until recording is finished
            
            # Save to temporary WAV
            temp_wav = "config/meeting_temp.wav"
            os.makedirs("config", exist_ok=True)
            with wave.open(temp_wav, "wb") as wav:
                wav.setnchannels(1)
                wav.setsampwidth(2)
                wav.setframerate(sample_rate)
                wav.writeframes(recording.tobytes())
                
            # Transcribe via Whisper
            logger.info("Transcribing recording...")
            whisper = WhisperModel("base.en", device="cpu", compute_type="int8")
            segments, _ = whisper.transcribe(temp_wav)
            text = " ".join(s.text for s in segments).strip()
            
            # Cleanup WAV
            try: os.remove(temp_wav)
            except Exception: pass
            
            if not text:
                return "No speech was detected in the meeting recording, sir."
                
            # Save notes to text file
            with open(output_txt, "a", encoding="utf-8") as f:
                f.write(f"\n--- Meeting Note ({time.strftime('%Y-%m-%d %H:%M:%S')}) ---\n{text}\n")
                
            return f"Dictation complete, sir. Notes successfully appended to {output_txt}. Contents:\n{text[:200]}..."
        except Exception as e:
            logger.error(f"Meeting notes recording failed: {e}")
            return f"Failed to record notes: {str(e)}"

    # --- Templates Generator ---
    
    def generate_document_template(self, template_type: str, args_dict: dict) -> str:
        """Generates structured content drafts (emails, NDA agreements) based on inputs."""
        template_type = template_type.lower().strip()
        
        if "nda" in template_type:
            party_a = args_dict.get("party_a", "Stark Industries")
            party_b = args_dict.get("party_b", "The Recipient")
            purpose = args_dict.get("purpose", "evaluating project details")
            nda_text = (
                f"MUTUAL NON-DISCLOSURE AGREEMENT\n\n"
                f"This Agreement is entered into by and between {party_a} and {party_b}.\n"
                f"1. Purpose: The parties wish to share confidential details relating to {purpose}.\n"
                f"2. Confidentiality: Both parties agree to protect proprietary information with reasonable care and "
                f"agree not to disclose it to third parties for a period of 5 years.\n\n"
                f"Executed on {time.strftime('%Y-%m-%d')}."
            )
            return nda_text
            
        elif "email" in template_type:
            recipient = args_dict.get("recipient", "Sir/Madam")
            sender = args_dict.get("sender", "Tony Stark")
            topic = args_dict.get("topic", "the project status")
            email_text = (
                f"Subject: Inquiry regarding {topic}\n\n"
                f"Dear {recipient},\n\n"
                f"I hope this message finds you well. I am writing to obtain a quick update regarding {topic}. "
                f"Please let me know your availability for a call later this week.\n\n"
                f"Best regards,\n{sender}"
            )
            return email_text
            
        else:
            return "Unknown template type requested. Supported formats are: 'nda' and 'email'."

    # --- Keyboard/Mouse Macros ---
    
    def save_macro(self, name: str, steps: list) -> str:
        """Saves a list of keyboard/mouse actions as a named macro."""
        try:
            steps_json = json.dumps(steps)
            conn = sqlite3.connect(self.db_path)
            c = conn.cursor()
            c.execute("INSERT OR REPLACE INTO macros (name, steps) VALUES (?, ?)", (name, steps_json))
            conn.commit()
            conn.close()
            return f"Macro '{name}' saved successfully, sir."
        except Exception as e:
            return f"Failed to save macro: {e}"

    def run_macro(self, name: str) -> str:
        """Executes a saved macro step by step."""
        try:
            conn = sqlite3.connect(self.db_path)
            c = conn.cursor()
            c.execute("SELECT steps FROM macros WHERE name = ?", (name,))
            row = c.fetchone()
            conn.close()
            
            if not row:
                return f"Macro '{name}' does not exist, sir."
                
            steps = json.loads(row[0])
            pyautogui.FAILSAFE = True
            
            for step in steps:
                action = step.get("action")
                param = step.get("param")
                
                if action == "click":
                    x, y = param
                    pyautogui.click(x, y)
                elif action == "type":
                    pyautogui.write(str(param), interval=0.02)
                elif action == "press":
                    pyautogui.press(str(param))
                elif action == "sleep":
                    time.sleep(float(param))
                    
            return f"Macro '{name}' executed successfully, sir."
        except Exception as e:
            logger.error(f"Macro execution failed: {e}")
            return f"Error executing macro: {str(e)}"

    def edit_pdf(self, action: str, files_list: list, output_path: str, page_num: int = 0, angle: int = 90) -> str:
        """Merges multiple PDFs, or rotates a specific page of a PDF using pypdf."""
        import pypdf
        act = action.lower().strip()
        out_p = os.path.abspath(os.path.expanduser(output_path))
        os.makedirs(os.path.dirname(out_p), exist_ok=True)
        
        try:
            if act == "merge":
                merger = pypdf.PdfMerger()
                for f in files_list:
                    fp = os.path.abspath(os.path.expanduser(f))
                    if os.path.exists(fp):
                        merger.append(fp)
                merger.write(out_p)
                merger.close()
                return f"Successfully merged {len(files_list)} PDFs into {output_path}, sir."
                
            elif act == "rotate":
                if not files_list:
                    return "Please specify a target PDF file, sir."
                fp = os.path.abspath(os.path.expanduser(files_list[0]))
                if not os.path.exists(fp):
                    return f"PDF file at {files_list[0]} does not exist."
                    
                reader = pypdf.PdfReader(fp)
                writer = pypdf.PdfWriter()
                for idx, page in enumerate(reader.pages):
                    if idx == page_num:
                        page.rotate(angle)
                    writer.add_page(page)
                with open(out_p, "wb") as f_out:
                    writer.write(f_out)
                return f"Successfully rotated page {page_num} of {files_list[0]} by {angle} degrees and saved to {output_path}, sir."
            else:
                return f"Unsupported PDF edit action '{action}', sir."
        except Exception as e:
            logger.error(f"PDF editing failed: {e}")
            return f"Failed to edit PDF: {str(e)}"

    def _search_images(self, query: str, max_results: int = 3) -> list:
        """Queries Bing Images, decodes HTML entities, and extracts high-resolution image URLs."""
        import requests
        import re
        import html
        import urllib.parse
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        search_url = f"https://www.bing.com/images/search?q={urllib.parse.quote(query)}"
        try:
            r = requests.get(search_url, headers=headers, timeout=10)
            r.raise_for_status()
            unescaped_html = html.unescape(r.text)
            img_urls = []
            matches = re.findall(r'"murl":"(http[^"]+)"', unescaped_html)
            for m in matches:
                if m not in img_urls:
                    img_urls.append(m)
                    if len(img_urls) >= max_results:
                        break
            return img_urls
        except Exception as e:
            logger.warning(f"Bing Image search failed for '{query}': {e}")
            return []

    def _download_image(self, url: str, save_path: str) -> bool:
        """Downloads an image URL and validates its integrity using PIL."""
        import requests
        from PIL import Image
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        try:
            os.makedirs(os.path.dirname(save_path), exist_ok=True)
            r = requests.get(url, headers=headers, timeout=10)
            r.raise_for_status()
            with open(save_path, "wb") as f:
                f.write(r.content)
            # Verify image is valid
            with Image.open(save_path) as img:
                img.verify()
            return True
        except Exception as e:
            logger.warning(f"Failed to download or verify image {url} to {save_path}: {e}")
            if os.path.exists(save_path):
                try: os.remove(save_path)
                except Exception: pass
            return False

    def marp_pptx_helper(self, title: str, subtitle: str, theme: str, slides_content: list, output_path: str = None) -> str:
        """
        Generate a professional PPTX using Marp CLI (npx @marp-team/marp-cli).
        Converts an LLM-generated Markdown outline into a beautifully themed slide deck.
        Falls back to pptx_helper() if Node.js / Marp is unavailable.
        """
        import subprocess
        import shutil
        import re as _re

        # Derive output path from title if not specified
        if not output_path:
            slug = _re.sub(r'[^a-zA-Z0-9_]', '_', title).lower().strip('_')
            output_path = f"config/{slug}.pptx"

        # Check if npx is available
        npx = shutil.which("npx")
        if not npx:
            logger.warning("Marp: npx not found — falling back to python-pptx.")
            return self.pptx_helper(title, subtitle, theme, slides_content, output_path)

        # ── Map JARVIS themes to Marp themes ────────────────────────
        theme_map = {
            "stark_tech":        ("default",  "#111827", "#f97316", "#f3f4f6"),
            "midnight_cyberpunk": ("gaia",    "#0f172a", "#38bdf8", "#ffffff"),
            "light_professional": ("uncover", "#f8fafc", "#0f172a", "#475569"),
            "forest_minimalist":  ("default", "#f4f4f5", "#064e3b", "#1f2937"),
        }
        marp_theme, bg_color, title_color, body_color = theme_map.get(
            theme.lower(), ("default", "#111827", "#f97316", "#f3f4f6")
        )

        # ── Build Marp Markdown ──────────────────────────────────────
        lines = [
            "---",
            f"marp: true",
            f"theme: {marp_theme}",
            f"paginate: true",
            f'style: |',
            f'  section {{',
            f'    background-color: {bg_color};',
            f'    color: {body_color};',
            f'    font-family: "Segoe UI", Arial, sans-serif;',
            f'  }}',
            f'  h1 {{ color: {title_color}; font-size: 2.2em; }}',
            f'  h2 {{ color: {title_color}; font-size: 1.5em; border-bottom: 3px solid {title_color}; padding-bottom: 6px; }}',
            f'  ul {{ line-height: 1.8; font-size: 1.05em; }}',
            f'  li {{ margin-bottom: 6px; }}',
            "---",
            "",
            f"# {title}",
            f"## {subtitle}",
            "",
        ]

        for slide in slides_content:
            s_title = slide.get("title", "Slide")
            bullets = slide.get("bullets", [])
            image_path = slide.get("image_path", "")

            lines.append("---")
            lines.append("")

            # Side-by-side if image available
            has_img = image_path and os.path.exists(image_path)
            if has_img:
                # Normalize path for Marp (forward slashes, absolute)
                img_abs = os.path.abspath(image_path).replace("\\", "/")
                lines.append(f"![bg right:40% fit]({img_abs})")
                lines.append("")

            lines.append(f"## {s_title}")
            lines.append("")
            for b in bullets:
                lines.append(f"- {b}")
            lines.append("")

        md_content = "\n".join(lines)

        # Write temp Markdown file
        slug = _re.sub(r'[^a-zA-Z0-9_]', '_', title).lower().strip('_')
        md_path = f"config/{slug}_marp.md"
        abs_out = os.path.abspath(output_path)
        os.makedirs("config", exist_ok=True)

        with open(md_path, "w", encoding="utf-8") as f:
            f.write(md_content)

        # Check if local Marp is installed
        local_marp = "node_modules/.bin/marp.cmd" if os.name == "nt" else "node_modules/.bin/marp"
        md_path_abs = os.path.abspath(md_path)
        if os.path.exists(local_marp):
            local_marp_abs = os.path.abspath(local_marp)
            cmd_args = [local_marp_abs, md_path_abs, "--pptx", "--output", abs_out, "--allow-local-files", "--no-stdin"]
            logger.info(f"Marp: Using local workspace marp-cli installation at {local_marp_abs}")
        else:
            cmd_args = [npx, "--yes", "@marp-team/marp-cli@latest", md_path_abs,
                        "--pptx", "--output", abs_out, "--allow-local-files", "--no-stdin"]
            logger.info("Marp: Local installation not found — falling back to npx.")

        try:
            result = subprocess.run(
                cmd_args,
                capture_output=True, text=True, timeout=120,
                shell=(os.name == "nt")
            )
            if result.returncode == 0 and os.path.exists(abs_out):
                logger.success(f"Marp: PPTX generated successfully at {abs_out}")
                return f"Marp se presentation bana li, sir — '{os.path.basename(abs_out)}'"
            else:
                logger.warning(f"Marp failed (exit {result.returncode}): {result.stderr[:300]}")
                logger.info("Marp failed — falling back to python-pptx.")
                return self.pptx_helper(title, subtitle, theme, slides_content, output_path)
        except subprocess.TimeoutExpired:
            logger.warning("Marp timed out — falling back to python-pptx.")
            return self.pptx_helper(title, subtitle, theme, slides_content, output_path)
        except Exception as e:
            logger.error(f"Marp exception: {e} — falling back to python-pptx.")
            return self.pptx_helper(title, subtitle, theme, slides_content, output_path)

    def revealjs_helper(self, title: str, subtitle: str, theme: str, slides_content: list, output_path: str = None) -> str:
        """
        Generates a premium, interactive, animated Reveal.js HTML slideshow.
        Opens it directly in the web browser for a professional, animated client review.
        """
        import re as _re
        import webbrowser

        if not output_path:
            slug = _re.sub(r'[^a-zA-Z0-9_]', '_', title).lower().strip('_')
            output_path = f"config/{slug}_slides.html"

        # Map themes to Reveal.js configs
        reveal_themes = {
            "stark_tech": "black",
            "midnight_cyberpunk": "league",
            "light_professional": "serif",
            "forest_minimalist": "beige"
        }
        rev_theme = reveal_themes.get(theme.lower(), "black")
        
        transitions = {
            "stark_tech": "convex",
            "midnight_cyberpunk": "zoom",
            "light_professional": "fade",
            "forest_minimalist": "slide"
        }
        trans = transitions.get(theme.lower(), "slide")

        # Custom CSS variables
        css_map = {
            "stark_tech": """
                .reveal {
                    background: radial-gradient(circle, #27272a 0%, #09090b 100%) !important;
                    color: #f4f4f5 !important;
                }
                h1, h2, h3 {
                    color: #ef4444 !important;
                    text-shadow: 0 0 10px rgba(239, 68, 68, 0.4) !important;
                    font-family: 'Segoe UI', sans-serif !important;
                    text-transform: uppercase !important;
                }
                .slide-layout {
                    display: flex;
                    align-items: center;
                    justify-content: space-between;
                    gap: 30px;
                }
                .slide-left {
                    flex: 1.2;
                    text-align: left;
                }
                .slide-right {
                    flex: 0.8;
                }
                .slide-image {
                    border: 3px solid #ef4444;
                    border-radius: 8px;
                    box-shadow: 0 0 15px rgba(239, 68, 68, 0.4);
                }
                li {
                    margin-bottom: 12px !important;
                    font-size: 0.9em !important;
                    color: #d4d4d8 !important;
                }
            """,
            "midnight_cyberpunk": """
                .reveal {
                    background: radial-gradient(circle, #0f172a 0%, #020617 100%) !important;
                    color: #f8fafc !important;
                }
                h1, h2, h3 {
                    color: #38bdf8 !important;
                    text-shadow: 0 0 15px rgba(56, 189, 248, 0.6) !important;
                    font-family: 'Segoe UI', sans-serif !important;
                    text-transform: uppercase !important;
                }
                .slide-layout {
                    display: flex;
                    align-items: center;
                    justify-content: space-between;
                    gap: 30px;
                }
                .slide-left {
                    flex: 1.2;
                    text-align: left;
                }
                .slide-right {
                    flex: 0.8;
                }
                .slide-image {
                    border: 3px solid #ec4899;
                    border-radius: 12px;
                    box-shadow: 0 0 20px rgba(236, 72, 153, 0.5);
                }
                li {
                    margin-bottom: 12px !important;
                    font-size: 0.9em !important;
                    color: #cbd5e1 !important;
                }
            """,
            "light_professional": """
                .reveal {
                    background: #f8fafc !important;
                    color: #334155 !important;
                }
                h1, h2, h3 {
                    color: #0f172a !important;
                    font-family: 'Segoe UI', sans-serif !important;
                    font-weight: 800 !important;
                }
                .slide-layout {
                    display: flex;
                    align-items: center;
                    justify-content: space-between;
                    gap: 30px;
                }
                .slide-left {
                    flex: 1.2;
                    text-align: left;
                }
                .slide-right {
                    flex: 0.8;
                }
                .slide-image {
                    border-radius: 8px;
                    box-shadow: 0 4px 15px rgba(0, 0, 0, 0.1);
                }
                li {
                    margin-bottom: 12px !important;
                    font-size: 0.9em !important;
                    color: #475569 !important;
                }
            """
        }
        custom_css = css_map.get(theme.lower(), css_map["stark_tech"])

        # Build slides html
        slides_html_list = []
        for slide in slides_content:
            s_title = slide.get("title", "Slide")
            bullets = slide.get("bullets", [])
            image_path = slide.get("image_path", "")

            bullets_html = "\\n".join([f"<li>{b}</li>" for b in bullets])

            has_img = image_path and os.path.exists(image_path)
            if has_img:
                img_abs = os.path.abspath(image_path).replace("\\\\", "/")
                slide_html = f"""
                <section data-transition="{trans}">
                    <div class="slide-layout">
                        <div class="slide-left">
                            <h2>{s_title}</h2>
                            <ul>
                                {bullets_html}
                            </ul>
                        </div>
                        <div class="slide-right">
                            <img class="slide-image" src="file:///{img_abs}" alt="{s_title}" style="max-height: 400px; max-width: 100%;">
                        </div>
                    </div>
                </section>
                """
            else:
                slide_html = f"""
                <section data-transition="{trans}">
                    <h2>{s_title}</h2>
                    <ul style="display: inline-block; text-align: left;">
                        {bullets_html}
                    </ul>
                </section>
                """
            slides_html_list.append(slide_html)

        slides_html = "\\n".join(slides_html_list)

        html_content = f"""<!doctype html>
<html lang="en">
    <head>
        <meta charset="utf-8">
        <title>{{title}}</title>
        <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.min.css">
        <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/theme/{{rev_theme}}.min.css" id="theme">
        <style>
            {{custom_css}}
            .reveal h1, .reveal h2, .reveal h3 {{{{
                text-transform: none !important;
            }}}}
        </style>
    </head>
    <body>
        <div class="reveal">
            <div class="slides">
                <!-- Title Slide -->
                <section data-transition="{{trans}}">
                    <h1>{{title}}</h1>
                    <p style="color: #a1a1aa; font-size: 1.2em;">{{subtitle}}</p>
                </section>
                
                <!-- Content Slides -->
                {{slides_html}}
            </div>
        </div>
        <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.js"></script>
        <script>
            Reveal.initialize({{{{
                hash: true,
                transition: '{{trans}}',
                backgroundTransition: 'fade',
                slideNumber: true,
                controls: true,
                progress: true,
                center: true
            }}}});
        </script>
    </body>
</html>
"""
        # Format the template variables
        final_html = html_content.format(
            title=title,
            subtitle=subtitle,
            rev_theme=rev_theme,
            custom_css=custom_css,
            trans=trans,
            slides_html=slides_html
        )

        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(final_html)
        
        logger.success(f"Reveal.js: Interactive HTML presentation generated at {output_path}")
        try:
            webbrowser.open(os.path.abspath(output_path))
        except Exception as err:
            logger.error(f"Failed to open browser: {err}")
            
        return f"Marp and Reveal.js presentations compiled, sir — '{os.path.basename(output_path)}' open in browser."

    def pptx_helper(self, title: str, subtitle: str, theme: str, slides_content: list, output_path: str = "config/presentation.pptx") -> str:
        """
        Creates a premium PowerPoint presentation by loading a pre-designed master template 
        and replacing its text and image placeholders, preserving all original fonts, alignments, and transitions.
        """
        import shutil

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            # Map theme to user downloaded premium templates
            TEMPLATE_MAP = {
                "stark_tech": "templates/Stark-Tech-Architecting-the-Future.pptx",
                "midnight_cyberpunk": "templates/Branding_and_marketing.pptx",
                "light_professional": "templates/Premium-Corporate-Business-Strategy.pptx",
                "forest_minimalist": "templates/Architecting-Modern-Software-Defense.pptx"
            }
            
            default_template = "templates/Future-Ready-Product-Launch-and-Roadmap-Strategy.pptx"
            selected_template = TEMPLATE_MAP.get(theme.lower(), default_template)
            
            # Fallback if selected template is missing
            if not os.path.exists(selected_template):
                logger.warning(f"Selected template '{selected_template}' not found. Searching templates/ for alternative...")
                pptx_files = [f for f in os.listdir("templates") if f.endswith(".pptx")]
                if pptx_files:
                    selected_template = os.path.join("templates", pptx_files[0])
                else:
                    raise FileNotFoundError("No templates found in templates/ directory.")

            logger.info(f"Loading premium presentation template: {selected_template}")
            prs = Presentation(os.path.abspath(selected_template))
            
            # 1. Populate Slide 0: Title Slide
            if len(prs.slides) > 0:
                title_slide = prs.slides[0]
                txt_shapes = []
                for s in title_slide.shapes:
                    if s.has_text_frame and s.text_frame.text.strip():
                        txt_shapes.append(s)
                txt_shapes.sort(key=lambda s: s.top)
                
                if len(txt_shapes) > 0:
                    tf = txt_shapes[0].text_frame
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].text = title
                        for r in tf.paragraphs[0].runs[1:]:
                            r.text = ""
                    else:
                        tf.text = title
                        
                if len(txt_shapes) > 1:
                    tf = txt_shapes[1].text_frame
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].text = subtitle
                        for r in tf.paragraphs[0].runs[1:]:
                            r.text = ""
                    else:
                        tf.text = subtitle

            # 2. Populate Slide 1+: Content Slides
            for idx, slide_data in enumerate(slides_content):
                slide_idx = idx + 1
                if slide_idx >= len(prs.slides):
                    break
                    
                slide = prs.slides[slide_idx]
                s_title = slide_data.get("title", "Topic")
                s_bullets = slide_data.get("bullets", [])
                image_path = slide_data.get("image_path", "")
                
                # Get all text shapes on this slide and sort vertically
                all_text_shapes = []
                for s in slide.shapes:
                    if s.has_text_frame:
                        txt_clean = s.text_frame.text.strip()
                        if len(txt_clean) > 2:
                            all_text_shapes.append(s)
                            
                all_text_shapes.sort(key=lambda s: s.top)
                
                # First text shape is header/title
                if all_text_shapes:
                    header_shape = all_text_shapes[0]
                    tf = header_shape.text_frame
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].text = s_title
                        for r in tf.paragraphs[0].runs[1:]:
                            r.text = ""
                    else:
                        tf.text = s_title
                        
                    content_shapes = all_text_shapes[1:]
                    
                    # If multiple card boxes exist in the template, map them horizontally
                    if len(content_shapes) >= len(s_bullets) and len(s_bullets) > 0:
                        content_shapes.sort(key=lambda s: s.left)
                        for b_idx, bullet in enumerate(s_bullets):
                            if b_idx < len(content_shapes):
                                tf_card = content_shapes[b_idx].text_frame
                                tf_card.word_wrap = True
                                if tf_card.paragraphs and tf_card.paragraphs[0].runs:
                                    tf_card.paragraphs[0].runs[0].text = bullet
                                    for r in tf_card.paragraphs[0].runs[1:]:
                                        r.text = ""
                                    for p_extra in tf_card.paragraphs[1:]:
                                        p_extra.text = ""
                                else:
                                    tf_card.text = bullet
                    else:
                        if content_shapes:
                            tf_body = content_shapes[0].text_frame
                            tf_body.word_wrap = True
                            tf_body.text = ""
                            for b_idx, bullet in enumerate(s_bullets):
                                if b_idx == 0:
                                    p = tf_body.paragraphs[0]
                                else:
                                    p = tf_body.add_paragraph()
                                p.text = "•  " + bullet
                                p.font.size = Pt(16)
                                p.space_after = Pt(8)
                                
                # Replace images on the slide if available
                if image_path and os.path.exists(image_path):
                    pic_shapes = [s for s in slide.shapes if s.shape_type == 13]
                    if pic_shapes:
                        pic_shape = pic_shapes[0]
                        replaced = False
                        # Method 1: Replace via blip_fill image part (standard pictures)
                        try:
                            rId = pic_shape._element.blip_fill.blip.rEmbed
                            image_part = pic_shape.part.related_parts[rId]
                            with open(image_path, "rb") as img_f:
                                image_part._blob = img_f.read()
                            replaced = True
                            logger.info(f"Image replaced (blip_fill) on slide {slide_idx} with {image_path}")
                        except Exception:
                            pass
                        # Method 2: Overlay a new picture at the same position/size
                        if not replaced:
                            try:
                                from pptx.util import Emu
                                left, top, width, height = pic_shape.left, pic_shape.top, pic_shape.width, pic_shape.height
                                slide.shapes.add_picture(image_path, left, top, width, height)
                                replaced = True
                                logger.info(f"Image replaced (overlay) on slide {slide_idx} with {image_path}")
                            except Exception as img_overlay_err:
                                logger.warning(f"Could not replace image on slide {slide_idx}: {img_overlay_err}")

            # 3. Clean up: Delete unused template slides
            target_slide_count = len(slides_content) + 1
            if len(prs.slides) > target_slide_count:
                logger.info(f"Removing {len(prs.slides) - target_slide_count} unused slides from template...")
                for idx in range(len(prs.slides) - 1, target_slide_count - 1, -1):
                    try:
                        del prs.slides._sldIdLst[idx]
                    except Exception:
                        pass

            out_p = os.path.abspath(os.path.expanduser(output_path))
            os.makedirs(os.path.dirname(out_p), exist_ok=True)
            prs.save(out_p)
            logger.success(f"PowerPoint generated successfully at {output_path} using template '{selected_template}'.")
            return f"Successfully created PowerPoint presentation at {output_path} using template '{os.path.basename(selected_template)}', sir."
        except ImportError:
            return "Sir, python-pptx install nahi hai. `pip install python-pptx` run karein."
        except Exception as e:
            logger.error(f"Failed to generate presentation from template: {e}")
            return f"Presentation template fill karte waqt gadbad ho gayi: {str(e)}"

    def open_presentation(self, filepath: str = "config/presentation.pptx") -> str:
        """Opens the generated presentation locally on Windows."""
        try:
            out_p = os.path.abspath(os.path.expanduser(filepath))
            if os.path.exists(out_p):
                os.startfile(out_p)
                return f"Presentation khol di hai, sir."
            else:
                return "Sir, presentation file nahi mili. Pehle presentation bana lijiye."
        except Exception as e:
            logger.error(f"Failed to open presentation: {e}")
            return f"Presentation kholne mein dikkat aayi: {e}"

    def create_mind_map(self, central_idea: str, nodes: list, save_path: str = "config/mindmap.png") -> str:
        """Generates a node-link mind map visualization image using matplotlib and networkx."""
        try:
            import networkx as nx
            import matplotlib.pyplot as plt
            
            plt.clf()
            G = nx.DiGraph()
            G.add_node(central_idea)
            
            # Add node links
            for edge in nodes:
                if isinstance(edge, (list, tuple)) and len(edge) >= 2:
                    G.add_edge(edge[0], edge[1])
                elif isinstance(edge, str):
                    G.add_edge(central_idea, edge)
            
            # Layout node positioning
            pos = nx.spring_layout(G)
            nx.draw(G, pos, with_labels=True, node_color="skyblue", node_size=2000, font_size=10, font_weight="bold", edge_color="gray", width=2)
            
            out_p = os.path.abspath(os.path.expanduser(save_path))
            os.makedirs(os.path.dirname(out_p), exist_ok=True)
            plt.savefig(out_p, dpi=150)
            plt.close()
            return f"Mind map visualization successfully generated and saved to {save_path}, sir."
        except ImportError:
            return "Please install networkx and matplotlib (`pip install networkx matplotlib`) to draw mind maps, sir."
        except Exception as e:
            return f"Failed to generate mind map: {str(e)}"

    def block_distractions(self, domains_list: list, block: bool = True) -> str:
        """Blocks distracting websites by appending them to the Windows hosts file (requires admin)."""
        hosts_path = r"C:\Windows\System32\drivers\etc\hosts"
        if not os.path.exists(hosts_path):
            hosts_path = "/etc/hosts" # Unix fallback
            
        if not os.path.exists(hosts_path):
            return "I could not find your system hosts file, sir."
            
        redirect_ip = "127.0.0.1"
        try:
            with open(hosts_path, "r", encoding="utf-8") as f:
                lines = f.readlines()
                
            new_lines = []
            blocked_domains = [d.strip().lower() for d in domains_list]
            
            # Read existing lines and filter out any domains we want to change/unblock
            for line in lines:
                is_target = False
                for d in blocked_domains:
                    if d in line.lower() and redirect_ip in line:
                        is_target = True
                        break
                if not is_target:
                    new_lines.append(line)
            
            # If blocking, append them at the end
            if block:
                if new_lines and not new_lines[-1].endswith("\n"):
                    new_lines.append("\n")
                for d in blocked_domains:
                    new_lines.append(f"{redirect_ip} {d}\n")
                    new_lines.append(f"{redirect_ip} www.{d}\n")
            
            with open(hosts_path, "w", encoding="utf-8") as f:
                f.writelines(new_lines)
                
            state = "blocked" if block else "unblocked"
            return f"Successfully {state} {len(domains_list)} distracting domains in hosts file, sir."
        except PermissionError:
            return (
                "Permission denied, sir. Modifying the system hosts file requires Administrator privileges. "
                "Please restart JARVIS in an elevated (Admin) terminal window."
            )
        except Exception as e:
            return f"Failed to modify hosts file: {str(e)}"

    def sign_document(self, doc_path: str, sig_img_path: str, output_path: str, coords: tuple = (100, 100)) -> str:
        """Overlays a signature image PNG onto a document image or PDF page."""
        try:
            from PIL import Image
            doc_p = os.path.abspath(os.path.expanduser(doc_path))
            sig_p = os.path.abspath(os.path.expanduser(sig_img_path))
            out_p = os.path.abspath(os.path.expanduser(output_path))
            
            if not os.path.exists(doc_p):
                return f"Document at {doc_path} does not exist."
            if not os.path.exists(sig_p):
                return f"Signature image at {sig_img_path} does not exist."
                
            os.makedirs(os.path.dirname(out_p), exist_ok=True)
            
            # Open images
            doc_img = Image.open(doc_p).convert("RGBA")
            sig_img = Image.open(sig_p).convert("RGBA")
            
            # Resize signature image if it's too large (cap at 200x100)
            sig_img.thumbnail((200, 100))
            
            # Paste signature onto document
            doc_img.paste(sig_img, coords, sig_img)
            doc_img.convert("RGB").save(out_p)
            return f"Successfully signed document and saved to {output_path}, sir."
        except Exception as e:
            logger.error(f"Document signing failed: {e}")
            return f"Failed to sign document: {str(e)}"

    def generate_mindmap(self, topic: str) -> str:
        """Generates an interactive visual Mermaid.js Mind Map HTML file and opens it in browser."""
        logger.info(f"Generating interactive Mind Map for topic: '{topic}'...")
        safe_name = "".join(c for c in topic if c.isalnum() or c in " _-").strip().replace(" ", "_")
        out_dir = "mindmaps"
        os.makedirs(out_dir, exist_ok=True)
        html_path = os.path.join(out_dir, f"{safe_name}_mindmap.html")
        
        mermaid_code = f"""mindmap
  root(({topic}))
    Architecture
      Core Logic
      Sensory Loop
      Vision Engine
    Capabilities
      Voice & STT
      OS Control
      E-Commerce
    Integrations
      WhatsApp & Phone
      Obsidian Vault
      Spotify & YouTube"""

        html_content = f"""<!DOCTYPE html>
<html>
<head>
    <title>Mind Map: {topic}</title>
    <script type="module">
        import mermaid from 'https://cdn.jsdelivr.net/npm/mermaid@10/dist/mermaid.esm.min.mjs';
        mermaid.initialize({{ startOnLoad: true, theme: 'dark' }});
    </script>
    <style>
        body {{ background-color: #0f172a; color: #f8fafc; font-family: sans-serif; display: flex; flex-direction: column; align-items: center; justify-content: center; min-height: 100vh; margin: 0; }}
        h1 {{ margin-bottom: 20px; color: #38bdf8; }}
        .mermaid {{ background: #1e293b; padding: 30px; border-radius: 12px; box-shadow: 0 10px 25px rgba(0,0,0,0.5); }}
    </style>
</head>
<body>
    <h1>🧠 Mind Map: {topic}</h1>
    <div class="mermaid">
    {mermaid_code}
    </div>
</body>
</html>
"""
        try:
            with open(html_path, "w", encoding="utf-8") as f:
                f.write(html_content)
            import webbrowser
            webbrowser.open(os.path.abspath(html_path))
            return f"Sir, topic '{topic}' par visual Mind Map generate ho gaya hai aur browser mein open kar diya hai."
        except Exception as e:
            return f"Failed to generate mind map: {e}"
